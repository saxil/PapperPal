import os
import requests
from bs4 import BeautifulSoup
from langchain.text_splitter import RecursiveCharacterTextSplitter
import fitz as PyMuPDF
import pdfplumber
from pptx import Presentation
from docx import Document

def load_and_chunk_file(file_path, chunk_size=1000, chunk_overlap=200):
    """Loads a file or URL, extracts text, and splits it into chunks."""
    if file_path.startswith("http"):
        text = extract_text_from_url(file_path)
        source = file_path
    else:
        file_extension = os.path.splitext(file_path)[1].lower()
        source = os.path.basename(file_path)

        if file_extension == ".pdf":
            text = extract_text_from_pdf(file_path)
        elif file_extension == ".pptx":
            text = extract_text_from_pptx(file_path)
        elif file_extension == ".docx":
            text = extract_text_from_docx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len
    )
    
    chunks = text_splitter.create_documents(
        texts=[text],
        metadatas=[{"source": source}]
    )

    return chunks

def extract_text_from_pdf(file_path):
    """Extracts text from a PDF file."""
    with pdfplumber.open(file_path) as pdf:
        text = "".join(page.extract_text() for page in pdf.pages if page.extract_text())
    return text

def extract_text_from_pptx(file_path):
    """Extracts text from a PowerPoint file."""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(file_path):
    """Extracts text from a Word document."""
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_url(url):
    """Extracts text from a URL."""
    try:
        response = requests.get(url)
        response.raise_for_status()  # Raise an exception for bad status codes
        soup = BeautifulSoup(response.content, "html.parser")
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return text
    except requests.exceptions.RequestException as e:
        return f"Error fetching URL: {e}"
